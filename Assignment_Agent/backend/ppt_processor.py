"""
PPT Processor - Extract text content from PowerPoint files
"""
import os
import base64
from pathlib import Path
from typing import Dict, List, Optional

# Optional import for PPTX processing
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    import logging
    logger = logging.getLogger(__name__)
    logger.warning("python-pptx library not available. Install with: pip install python-pptx")

# Optional import for legacy PPT files (requires comtypes on Windows)
try:
    import comtypes.client
    COMTYPES_AVAILABLE = True
except ImportError:
    COMTYPES_AVAILABLE = False


class PPTProcessor:
    """Process PowerPoint files and extract text content"""
    
    @staticmethod
    def is_ppt_file(file_path: str) -> bool:
        """Check if file is a PowerPoint file"""
        ext = Path(file_path).suffix.lower()
        return ext in ['.ppt', '.pptx', '.pptm']
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> Dict[str, any]:
        """
        Extract text from PPTX/PPTM file
        Returns dict with slides_text, total_slides, and slide_details
        """
        if not PPTX_AVAILABLE:
            import logging
            logger = logging.getLogger(__name__)
            logger.error("python-pptx library not available. Cannot process PPTX files.")
            return {
                'slides_text': '[python-pptx library not available. Install with: pip install python-pptx]',
                'total_slides': 0,
                'slide_details': []
            }
        
        try:
            import logging
            logger = logging.getLogger(__name__)
            logger.info(f"Attempting to extract text from PPTX file: {file_path}")
            
            prs = Presentation(file_path)
            slides_text = []
            slide_details = []
            
            logger.info(f"Found {len(prs.slides)} slides in presentation")
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text_parts = []
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    # Try direct text attribute first
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            slide_text_parts.append(text)
                    
                    # Also try text_frame (for text boxes and placeholders)
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        try:
                            # Get text from text_frame
                            frame_text = shape.text_frame.text.strip()
                            if frame_text and frame_text not in slide_text_parts:
                                slide_text_parts.append(frame_text)
                            
                            # Also check paragraphs in text_frame
                            for paragraph in shape.text_frame.paragraphs:
                                para_text = paragraph.text.strip()
                                if para_text and para_text not in slide_text_parts:
                                    slide_text_parts.append(para_text)
                        except Exception:
                            pass
                    
                    # Also check for tables
                    if hasattr(shape, "has_table") and shape.has_table:
                        try:
                            table_text = []
                            for row in shape.table.rows:
                                row_text = []
                                for cell in row.cells:
                                    cell_text = cell.text.strip() if cell.text else ""
                                    if cell_text:
                                        row_text.append(cell_text)
                                if row_text:
                                    table_text.append(" | ".join(row_text))
                            if table_text:
                                slide_text_parts.append("Table:\n" + "\n".join(table_text))
                        except Exception:
                            pass
                
                slide_text = "\n".join(slide_text_parts)
                
                if slide_text.strip():
                    slide_details.append({
                        'slide_number': slide_num,
                        'text': slide_text
                    })
                    slides_text.append(f"--- Slide {slide_num} ---\n{slide_text}")
            
            total_slides = len(prs.slides)
            combined_text = "\n\n".join(slides_text) if slides_text else "[No text content found in slides]"
            
            logger.info(f"Extracted text from {len(slide_details)} slides, total length: {len(combined_text)}")
            
            return {
                'slides_text': combined_text,
                'total_slides': total_slides,
                'slide_details': slide_details
            }
            
        except Exception as e:
            import logging
            logger = logging.getLogger(__name__)
            logger.error(f"Error reading PPTX file {file_path}: {e}", exc_info=True)
            return {
                'slides_text': f'[Error reading PPTX file: {str(e)}]',
                'total_slides': 0,
                'slide_details': []
            }
    
    @staticmethod
    def extract_text_from_ppt(file_path: str) -> Dict[str, any]:
        """
        Extract text from legacy PPT file using COM automation (Windows only)
        Returns dict with slides_text, total_slides, and slide_details
        """
        if not COMTYPES_AVAILABLE:
            return {
                'slides_text': '[comtypes library not available. Legacy PPT files require Windows and comtypes. Install with: pip install comtypes]',
                'total_slides': 0,
                'slide_details': []
            }
        
        try:
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = 1
            
            try:
                presentation = powerpoint.Presentations.Open(str(Path(file_path).absolute()))
                slides_text = []
                slide_details = []
                
                total_slides = presentation.Slides.Count
                
                for slide_num in range(1, total_slides + 1):
                    slide = presentation.Slides(slide_num)
                    slide_text_parts = []
                    
                    # Extract text from all shapes
                    for shape_num in range(1, slide.Shapes.Count + 1):
                        shape = slide.Shapes(shape_num)
                        if hasattr(shape, "TextFrame") and shape.TextFrame:
                            if hasattr(shape.TextFrame, "TextRange") and shape.TextFrame.TextRange:
                                text = shape.TextFrame.TextRange.Text.strip()
                                if text:
                                    slide_text_parts.append(text)
                    
                    slide_text = "\n".join(slide_text_parts)
                    
                    if slide_text.strip():
                        slide_details.append({
                            'slide_number': slide_num,
                            'text': slide_text
                        })
                        slides_text.append(f"--- Slide {slide_num} ---\n{slide_text}")
                
                combined_text = "\n\n".join(slides_text) if slides_text else "[No text content found in slides]"
                
                presentation.Close()
                powerpoint.Quit()
                
                return {
                    'slides_text': combined_text,
                    'total_slides': total_slides,
                    'slide_details': slide_details
                }
                
            except Exception as e:
                try:
                    powerpoint.Quit()
                except:
                    pass
                return {
                    'slides_text': f'[Error reading PPT file: {str(e)}]',
                    'total_slides': 0,
                    'slide_details': []
                }
                
        except Exception as e:
            return {
                'slides_text': f'[Error opening PowerPoint application: {str(e)}]',
                'total_slides': 0,
                'slide_details': []
            }
    
    @staticmethod
    def process_ppt_file(file_path: str) -> Dict[str, any]:
        """
        Process a PowerPoint file and extract text content
        Returns dict with slides_text, total_slides, and slide_details
        """
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pptx' or ext == '.pptm':
            return PPTProcessor.extract_text_from_pptx(file_path)
        elif ext == '.ppt':
            return PPTProcessor.extract_text_from_ppt(file_path)
        else:
            return {
                'slides_text': f'[Unsupported PowerPoint format: {ext}]',
                'total_slides': 0,
                'slide_details': []
            }
    
    @staticmethod
    def convert_slides_to_images_pptx(file_path: str) -> List[str]:
        """
        Convert PPTX slides to base64-encoded PNG images
        Returns list of base64 strings (one per slide)
        """
        if not PPTX_AVAILABLE:
            import logging
            logger = logging.getLogger(__name__)
            logger.error("python-pptx library not available. Cannot convert slides to images.")
            return []
        
        try:
            import logging
            logger = logging.getLogger(__name__)
            logger.info(f"Converting PPTX slides to images: {file_path}")
            
            prs = Presentation(file_path)
            slide_images = []
            
            # Try to use COM automation on Windows to export slides as images
            if COMTYPES_AVAILABLE:
                try:
                    import tempfile
                    import os
                    temp_dir = tempfile.mkdtemp()
                    
                    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
                    powerpoint.Visible = 0  # Don't show PowerPoint
                    
                    try:
                        presentation = powerpoint.Presentations.Open(str(Path(file_path).absolute()))
                        
                        for slide_num in range(1, len(prs.slides) + 1):
                            slide = presentation.Slides(slide_num)
                            image_path = os.path.join(temp_dir, f"slide_{slide_num}.png")
                            
                            # Export slide as PNG
                            slide.Export(image_path, "PNG", 1920, 1080)  # High resolution
                            
                            # Read image and convert to base64
                            with open(image_path, "rb") as img_file:
                                img_data = img_file.read()
                                img_base64 = base64.b64encode(img_data).decode('utf-8')
                                slide_images.append(img_base64)
                        
                        presentation.Close()
                        powerpoint.Quit()
                        
                        # Cleanup temp directory
                        import shutil
                        shutil.rmtree(temp_dir, ignore_errors=True)
                        
                        logger.info(f"Successfully converted {len(slide_images)} slides to images")
                        return slide_images
                        
                    except Exception as e:
                        try:
                            powerpoint.Quit()
                        except:
                            pass
                        import shutil
                        shutil.rmtree(temp_dir, ignore_errors=True)
                        logger.warning(f"COM automation failed, trying alternative method: {e}")
                except Exception as e:
                    logger.warning(f"COM automation not available or failed: {e}")
            
            # Fallback: Try using python-pptx with PIL (limited - won't render properly)
            # This is a placeholder - actual rendering requires Office automation or conversion service
            logger.warning("Direct slide rendering not available. COM automation required for image conversion.")
            return []
            
        except Exception as e:
            import logging
            logger = logging.getLogger(__name__)
            logger.error(f"Error converting slides to images: {e}", exc_info=True)
            return []
    
    @staticmethod
    def convert_slides_to_images(file_path: str) -> List[str]:
        """
        Convert PowerPoint slides to base64-encoded PNG images
        Returns list of base64 strings (one per slide)
        """
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pptx' or ext == '.pptm':
            return PPTProcessor.convert_slides_to_images_pptx(file_path)
        elif ext == '.ppt':
            # For legacy PPT, use COM automation
            if COMTYPES_AVAILABLE:
                try:
                    import logging
                    import tempfile
                    import os
                    import base64
                    logger = logging.getLogger(__name__)
                    
                    temp_dir = tempfile.mkdtemp()
                    slide_images = []
                    
                    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
                    powerpoint.Visible = 0
                    
                    try:
                        presentation = powerpoint.Presentations.Open(str(Path(file_path).absolute()))
                        total_slides = presentation.Slides.Count
                        
                        for slide_num in range(1, total_slides + 1):
                            slide = presentation.Slides(slide_num)
                            image_path = os.path.join(temp_dir, f"slide_{slide_num}.png")
                            
                            slide.Export(image_path, "PNG", 1920, 1080)
                            
                            with open(image_path, "rb") as img_file:
                                img_data = img_file.read()
                                img_base64 = base64.b64encode(img_data).decode('utf-8')
                                slide_images.append(img_base64)
                        
                        presentation.Close()
                        powerpoint.Quit()
                        
                        import shutil
                        shutil.rmtree(temp_dir, ignore_errors=True)
                        
                        return slide_images
                    except Exception as e:
                        try:
                            powerpoint.Quit()
                        except:
                            pass
                        import shutil
                        shutil.rmtree(temp_dir, ignore_errors=True)
                        logger.error(f"Error converting PPT slides: {e}")
                        return []
                except Exception as e:
                    import logging
                    logger = logging.getLogger(__name__)
                    logger.error(f"Error with COM automation: {e}")
                    return []
            else:
                return []
        else:
            return []
    
    @staticmethod
    def extract_design_metadata_pptx(file_path: str) -> Dict[str, any]:
        """
        Extract design metadata from PPTX file without requiring PowerPoint
        Returns dict with design_description, total_slides, and design_details
        """
        if not PPTX_AVAILABLE:
            import logging
            logger = logging.getLogger(__name__)
            logger.error("python-pptx library not available. Cannot extract design metadata.")
            return {
                'design_description': '[python-pptx library not available. Install with: pip install python-pptx]',
                'total_slides': 0,
                'design_details': []
            }
        
        try:
            import logging
            logger = logging.getLogger(__name__)
            logger.info(f"Extracting design metadata from PPTX file: {file_path}")
            
            prs = Presentation(file_path)
            design_details = []
            design_parts = []
            
            # Extract theme colors if available
            theme_colors = []
            try:
                if hasattr(prs, 'slide_master') and prs.slide_master:
                    if hasattr(prs.slide_master, 'theme') and prs.slide_master.theme:
                        if hasattr(prs.slide_master.theme, 'theme_part') and prs.slide_master.theme.theme_part:
                            # Try to get color scheme
                            pass
            except Exception:
                pass
            
            # Analyze each slide
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_design = {
                    'slide_number': slide_num,
                    'background': {},
                    'shapes': [],
                    'fonts': [],
                    'colors': [],
                    'layout_info': {}
                }
                
                # Extract background information
                try:
                    if hasattr(slide, 'background') and slide.background:
                        if hasattr(slide.background, 'fill'):
                            fill = slide.background.fill
                            bg_info = {}
                            if hasattr(fill, 'type'):
                                bg_info['type'] = str(fill.type)
                            if hasattr(fill, 'fore_color') and fill.fore_color:
                                if hasattr(fill.fore_color, 'rgb'):
                                    bg_info['color'] = str(fill.fore_color.rgb)
                            if bg_info:
                                slide_design['background'] = bg_info
                except Exception:
                    pass
                
                # Extract information from shapes
                shape_count = 0
                text_shapes = 0
                image_shapes = 0
                table_shapes = 0
                auto_shapes = 0
                
                for shape in slide.shapes:
                    shape_count += 1
                    shape_info = {
                        'type': type(shape).__name__,
                        'position': {},
                        'size': {},
                        'formatting': {}
                    }
                    
                    # Get position and size
                    try:
                        if hasattr(shape, 'left'):
                            shape_info['position']['left'] = shape.left
                        if hasattr(shape, 'top'):
                            shape_info['position']['top'] = shape.top
                        if hasattr(shape, 'width'):
                            shape_info['size']['width'] = shape.width
                        if hasattr(shape, 'height'):
                            shape_info['size']['height'] = shape.height
                    except Exception:
                        pass
                    
                    # Check shape type
                    try:
                        if hasattr(shape, 'image'):
                            image_shapes += 1
                            shape_info['has_image'] = True
                        elif hasattr(shape, 'has_table') and shape.has_table:
                            table_shapes += 1
                            shape_info['has_table'] = True
                        elif hasattr(shape, 'auto_shape_type'):
                            # Try to access auto_shape_type - it raises ValueError if not an auto shape
                            try:
                                auto_type = shape.auto_shape_type
                                auto_shapes += 1
                                shape_info['auto_shape_type'] = str(auto_type)
                            except ValueError:
                                # Not an auto shape, skip
                                pass
                    except Exception:
                        # Skip shape type detection if there's any error
                        pass
                    
                    # Extract text formatting
                    try:
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            text_shapes += 1
                            # Get font information from paragraphs
                            for paragraph in shape.text_frame.paragraphs:
                                if hasattr(paragraph, 'font') and paragraph.font:
                                    font_info = {}
                                    if hasattr(paragraph.font, 'name'):
                                        font_info['name'] = paragraph.font.name
                                    if hasattr(paragraph.font, 'size'):
                                        font_info['size'] = paragraph.font.size
                                    if hasattr(paragraph.font, 'bold'):
                                        font_info['bold'] = paragraph.font.bold
                                    if hasattr(paragraph.font, 'italic'):
                                        font_info['italic'] = paragraph.font.italic
                                    if hasattr(paragraph.font, 'color') and paragraph.font.color:
                                        if hasattr(paragraph.font.color, 'rgb'):
                                            font_info['color'] = str(paragraph.font.color.rgb)
                                    if font_info:
                                        shape_info['formatting']['font'] = font_info
                                        if font_info.get('name'):
                                            if font_info['name'] not in [f.get('name') for f in slide_design['fonts']]:
                                                slide_design['fonts'].append(font_info)
                                        if font_info.get('color'):
                                            if font_info['color'] not in slide_design['colors']:
                                                slide_design['colors'].append(font_info['color'])
                    except Exception:
                        pass
                    
                    # Extract fill colors
                    try:
                        if hasattr(shape, 'fill'):
                            fill = shape.fill
                            if hasattr(fill, 'fore_color') and fill.fore_color:
                                if hasattr(fill.fore_color, 'rgb'):
                                    color = str(fill.fore_color.rgb)
                                    if color not in slide_design['colors']:
                                        slide_design['colors'].append(color)
                    except Exception:
                        pass
                    
                    # Extract line colors
                    try:
                        if hasattr(shape, 'line'):
                            line = shape.line
                            if hasattr(line, 'color') and line.color:
                                if hasattr(line.color, 'rgb'):
                                    color = str(line.color.rgb)
                                    if color not in slide_design['colors']:
                                        slide_design['colors'].append(color)
                    except Exception:
                        pass
                    
                    slide_design['shapes'].append(shape_info)
                
                # Layout information
                slide_design['layout_info'] = {
                    'total_shapes': shape_count,
                    'text_shapes': text_shapes,
                    'image_shapes': image_shapes,
                    'table_shapes': table_shapes,
                    'auto_shapes': auto_shapes
                }
                
                design_details.append(slide_design)
                
                # Build text description for this slide
                slide_desc = [f"=== Slide {slide_num} Design ==="]
                
                if slide_design['background']:
                    slide_desc.append(f"Background: {slide_design['background']}")
                
                slide_desc.append(f"Layout: {shape_count} shapes ({text_shapes} text, {image_shapes} images, {table_shapes} tables, {auto_shapes} auto-shapes)")
                
                if slide_design['fonts']:
                    unique_fonts = {}
                    for font in slide_design['fonts']:
                        name = font.get('name', 'Unknown')
                        if name not in unique_fonts:
                            unique_fonts[name] = font
                    font_list = [f"{f.get('name')} (size: {f.get('size')}, bold: {f.get('bold')}, italic: {f.get('italic')})" 
                                for f in unique_fonts.values()]
                    slide_desc.append(f"Typography: {', '.join(font_list)}")
                
                if slide_design['colors']:
                    slide_desc.append(f"Colors used: {', '.join(slide_design['colors'])}")
                
                # Shape positioning analysis
                if slide_design['shapes']:
                    positions = [s.get('position', {}) for s in slide_design['shapes']]
                    sizes = [s.get('size', {}) for s in slide_design['shapes']]
                    slide_desc.append(f"Shape positioning: {len(positions)} positioned elements")
                
                design_parts.append("\n".join(slide_desc))
            
            total_slides = len(prs.slides)
            combined_description = "\n\n".join(design_parts) if design_parts else "[No design information found]"
            
            logger.info(f"Extracted design metadata from {len(design_details)} slides")
            
            return {
                'design_description': combined_description,
                'total_slides': total_slides,
                'design_details': design_details
            }
            
        except Exception as e:
            import logging
            logger = logging.getLogger(__name__)
            logger.error(f"Error extracting design metadata from PPTX file {file_path}: {e}", exc_info=True)
            return {
                'design_description': f'[Error extracting design metadata: {str(e)}]',
                'total_slides': 0,
                'design_details': []
            }
    
    @staticmethod
    def extract_design_metadata(file_path: str) -> Dict[str, any]:
        """
        Extract design metadata from PowerPoint file
        Returns dict with design_description, total_slides, and design_details
        """
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pptx' or ext == '.pptm':
            return PPTProcessor.extract_design_metadata_pptx(file_path)
        elif ext == '.ppt':
            # For legacy PPT files, we can't extract design metadata without PowerPoint
            return {
                'design_description': '[Legacy PPT format (.ppt) requires PowerPoint installation for design metadata extraction. Please use PPTX format for design evaluation without PowerPoint.]',
                'total_slides': 0,
                'design_details': []
            }
        else:
            return {
                'design_description': f'[Unsupported PowerPoint format: {ext}]',
                'total_slides': 0,
                'design_details': []
            }
    
    @staticmethod
    def process_multiple_ppt_files(file_paths: List[str]) -> List[Dict[str, any]]:
        """
        Process multiple PPT files
        Returns list of dicts, each containing slides_text, total_slides, and slide_details
        """
        results = []
        for file_path in file_paths:
            result = PPTProcessor.process_ppt_file(file_path)
            result['file_path'] = file_path
            result['filename'] = Path(file_path).name
            results.append(result)
        return results

